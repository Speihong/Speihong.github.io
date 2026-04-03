from pptx import Presentation
from pptx.util import Inches
import os

# 创建PPT
prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9 宽度
prs.slide_height = Inches(7.5)     # 16:9 高度

slides_dir = 'slides_v3'

# 获取所有slide图片并排序
slides = sorted([f for f in os.listdir(slides_dir) if f.endswith('.jpg')])

for slide_file in slides:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    
    img_path = os.path.join(slides_dir, slide_file)
    
    # 计算图片尺寸以适应幻灯片
    left = Inches(0)
    top = Inches(0)
    width = Inches(13.333)
    height = Inches(7.5)
    
    slide.shapes.add_picture(img_path, left, top, width, height)

# 保存PPT
prs.save('转正汇报_v3.pptx')
print(f'PPT已生成: 转正汇报_v3.pptx')
print(f'共 {len(slides)} 页')
